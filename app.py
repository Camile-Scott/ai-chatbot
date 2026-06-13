from flask import Flask, request, jsonify, render_template
from openai import OpenAI
import PyPDF2
from docx import Document
from pptx import Presentation
import sqlite3
import json
import re
import requests

app = Flask(__name__)

client = OpenAI(
    api_key="sk-b4ea12a9c50e4704a15f2a7ad5010167",
    base_url="https://dashscope.aliyuncs.com/compatible-mode/v1"
)


def init_db():
    conn = sqlite3.connect("chat.db")
    cursor = conn.cursor()
    cursor.execute("""
                   CREATE TABLE IF NOT EXISTS conversations
                   (
                       id
                       INTEGER
                       PRIMARY
                       KEY
                       AUTOINCREMENT,
                       title
                       TEXT
                       DEFAULT
                       '新对话',
                       created_at
                       TEXT
                       DEFAULT
                       CURRENT_TIMESTAMP
                   )
                   """)
    cursor.execute("""
                   CREATE TABLE IF NOT EXISTS messages
                   (
                       id
                       INTEGER
                       PRIMARY
                       KEY
                       AUTOINCREMENT,
                       conversation_id
                       INTEGER,
                       role
                       TEXT,
                       content
                       TEXT,
                       created_at
                       TEXT
                       DEFAULT
                       CURRENT_TIMESTAMP
                   )
                   """)
    cursor.execute("""
                   CREATE TABLE IF NOT EXISTS documents
                   (
                       id
                       INTEGER
                       PRIMARY
                       KEY
                       AUTOINCREMENT,
                       conversation_id
                       INTEGER,
                       name
                       TEXT,
                       content
                       TEXT
                   )
                   """)
    conn.commit()
    conn.close()


init_db()


def read_pdf(file):
    reader = PyPDF2.PdfReader(file)
    text = ""
    for page in reader.pages:
        text += page.extract_text()
    return text


def read_word(file):
    doc = Document(file)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text


def read_pptx(file):
    prs = Presentation(file)
    text = ""
    for i, slide in enumerate(prs.slides):
        text += f"第{i + 1}页:\n"
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text += para.text + "\n"
        text += "\n"
    return text


def calculate(expression):
    """计算数学表达式"""
    try:
        result = eval(expression)
        return str(result)
    except:
        return "计算错误"


def search_web(query):
    """网页搜索"""
    try:
        # 使用DuckDuckGo搜索
        search_url = f"https://duckduckgo.com/api?q={query}&format=json"
        response = requests.get(search_url, timeout=5)
        data = response.json()

        results = data.get("Results", [])
        if not results:
            return f"搜索'{query}'无结果"

        # 返回前3个结果
        summary = "\n".join([f"• {r.get('Title', '')}: {r.get('Text', '')[:100]}" for r in results[:3]])
        return summary if summary else "无搜索结果"
    except Exception as e:
        return f"搜索失败: {str(e)}"


@app.route("/")
def index():
    return render_template("index.html")


@app.route("/conversations", methods=["POST"])
def create_conversation():
    conn = sqlite3.connect("chat.db")
    cursor = conn.cursor()
    cursor.execute("INSERT INTO conversations (title) VALUES ('新对话')")
    conn.commit()
    new_id = cursor.lastrowid
    conn.close()
    return jsonify({"id": new_id, "title": "新对话"})


@app.route("/conversations", methods=["GET"])
def get_conversations():
    conn = sqlite3.connect("chat.db")
    cursor = conn.cursor()
    cursor.execute("SELECT id, title FROM conversations ORDER BY id DESC")
    rows = cursor.fetchall()
    conn.close()
    return jsonify([{"id": r[0], "title": r[1]} for r in rows])


@app.route("/conversations/<int:conv_id>", methods=["DELETE"])
def delete_conversation(conv_id):
    conn = sqlite3.connect("chat.db")
    cursor = conn.cursor()
    cursor.execute("DELETE FROM conversations WHERE id = ?", (conv_id,))
    cursor.execute("DELETE FROM messages WHERE conversation_id = ?", (conv_id,))
    cursor.execute("DELETE FROM documents WHERE conversation_id = ?", (conv_id,))
    conn.commit()
    conn.close()
    return jsonify({"success": True})


@app.route("/chat", methods=["POST"])
def chat():
    user_input = request.json.get("message")
    conversation_id = request.json.get("conversation_id")

    conn = sqlite3.connect("chat.db")
    cursor = conn.cursor()

    cursor.execute("SELECT role, content FROM messages WHERE conversation_id = ? ORDER BY id", (conversation_id,))
    history = [{"role": r[0], "content": r[1]} for r in cursor.fetchall()]
    history.append({"role": "user", "content": user_input})

    # 添加系统提示，告诉AI什么时候用工具
    system_prompt = """你是一个有用的助手。你有两个工具可以使用：

1. 计算器：当用户问数学计算问题时使用
   使用格式：[CALC]表达式[/CALC]，例如：[CALC]100+200[/CALC]

2. 网页搜索：当用户问最新信息、新闻、或需要查资料时使用
   使用格式：[SEARCH]搜索内容[/SEARCH]，例如：[SEARCH]今天的新闻[/SEARCH]

如果用户问的不是这些问题，就直接回答。"""

    response = client.chat.completions.create(
        model="deepseek-v3",
        messages=[{"role": "system", "content": system_prompt}] + history
    )
    reply = response.choices[0].message.content

    # 检查回复中是否有计算标签，如果有则执行计算
    calc_pattern = r'\[CALC\](.*?)\[/CALC\]'
    matches = re.findall(calc_pattern, reply)

    for match in matches:
        result = calculate(match)
        reply = reply.replace(f'[CALC]{match}[/CALC]', f'{match}={result}')

    # 检查回复中是否有搜索标签，如果有则执行网页搜索
    search_pattern = r'\[SEARCH\](.*?)\[/SEARCH\]'
    search_matches = re.findall(search_pattern, reply)

    for match in search_matches:
        result = search_web(match)
        reply = reply.replace(f'[SEARCH]{match}[/SEARCH]', f'搜索"{match}"的结果:\n{result}')

    cursor.execute("INSERT INTO messages (conversation_id, role, content) VALUES (?, ?, ?)",
                   (conversation_id, "user", user_input))
    cursor.execute("INSERT INTO messages (conversation_id, role, content) VALUES (?, ?, ?)",
                   (conversation_id, "assistant", reply))

    cursor.execute("SELECT COUNT(*) FROM messages WHERE conversation_id = ?", (conversation_id,))
    count = cursor.fetchone()[0]
    if count == 2:
        title = user_input[:20]
        cursor.execute("UPDATE conversations SET title = ? WHERE id = ?", (title, conversation_id))

    conn.commit()
    conn.close()

    return jsonify({"reply": reply})


@app.route("/history", methods=["GET"])
def history():
    conversation_id = request.args.get("conversation_id")
    conn = sqlite3.connect("chat.db")
    cursor = conn.cursor()
    cursor.execute("SELECT role, content FROM messages WHERE conversation_id = ? ORDER BY id", (conversation_id,))
    rows = cursor.fetchall()
    conn.close()
    return jsonify([{"role": r[0], "content": r[1]} for r in rows])


@app.route("/upload", methods=["POST"])
def upload():
    file = request.files["file"]
    conversation_id = request.form.get("conversation_id")
    filename = file.filename.lower()

    if filename.endswith(".pdf"):
        text = read_pdf(file)
    elif filename.endswith(".docx"):
        text = read_word(file)
    elif filename.endswith(".pptx"):
        text = read_pptx(file)
    else:
        return jsonify({"message": "不支持的文件格式"}), 400

    conn = sqlite3.connect("chat.db")
    cursor = conn.cursor()
    cursor.execute("INSERT INTO documents (conversation_id, name, content) VALUES (?, ?, ?)",
                   (conversation_id, file.filename, text))
    conn.commit()

    cursor.execute("SELECT name FROM documents WHERE conversation_id = ?", (conversation_id,))
    names = [r[0] for r in cursor.fetchall()]
    conn.close()

    return jsonify({
        "message": f"上传成功!当前已加载 {len(names)} 份文档:{', '.join(names)}"
    })


@app.route("/ask", methods=["POST"])
def ask():
    question = request.json.get("question")
    conversation_id = request.json.get("conversation_id")

    conn = sqlite3.connect("chat.db")
    cursor = conn.cursor()
    cursor.execute("SELECT name, content FROM documents WHERE conversation_id = ?", (conversation_id,))
    docs = cursor.fetchall()
    conn.close()

    combined = "\n\n".join([f"【{d[0]}】\n{d[1]}" for d in docs])
    response = client.chat.completions.create(
        model="deepseek-v3",
        messages=[
            {"role": "system", "content": f"你是文档助手。以下是所有文档内容:\n\n{combined}"},
            {"role": "user", "content": question}
        ]
    )
    return jsonify({"answer": response.choices[0].message.content})


if __name__ == "__main__":
    app.run(debug=True)